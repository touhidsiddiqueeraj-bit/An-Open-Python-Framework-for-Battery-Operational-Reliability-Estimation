#!/usr/bin/env python3
"""Convert presentation.html content into a .pptx file."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x60, 0xA5, 0xFA)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
RED = RGBColor(0xF8, 0x71, 0x71)
GREEN = RGBColor(0x34, 0xD3, 0x99)
WHITE = RGBColor(0xE2, 0xE8, 0xF0)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
DIM = RGBColor(0x64, 0x7B, 0x48)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)


def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_frame(slide, left, top, width, height, items, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return tf


def add_card(slide, left, top, width, height, title, body, title_color=RED):
    """Add a rounded-rect card with title + body."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()
    # Title
    add_text_box(slide, left + 0.3, top + 0.15, width - 0.6, 0.5,
                 title, font_size=15, bold=True, color=title_color)
    # Body
    add_text_box(slide, left + 0.3, top + 0.55, width - 0.6, height - 0.7,
                 body, font_size=13, color=WHITE)


def add_stat_box(slide, left, top, num, label):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(2.5), Inches(1.8)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()
    add_text_box(slide, left, top + 0.15, 2.5, 0.8,
                 str(num), font_size=36, bold=True, color=ACCENT,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + 0.95, 2.5, 0.6,
                 label, font_size=12, color=GRAY,
                 alignment=PP_ALIGN.CENTER)


def add_table(slide, left, top, width, headers, rows, col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                          Inches(left), Inches(top),
                                          Inches(width), Inches(0.4 * n_rows))
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            p.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = WHITE
                p.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_DARK

    return table


# ── SLIDE 1: Title ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)
add_text_box(slide, 0.8, 1.2, 11.5, 1.2,
             "An Open Python Framework for\nBattery Operational Reliability Estimation",
             font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
add_text_box(slide, 0.8, 2.6, 11.5, 0.6,
             "With Quantitative Minimum-Data Guidelines from a Synthetic Scaling Study",
             font_size=22, color=ACCENT, alignment=PP_ALIGN.LEFT)
add_text_box(slide, 0.8, 3.3, 11.5, 0.4,
             "Team Dynamic  •  June 2026",
             font_size=16, color=GRAY)

for i, (num, label) in enumerate([("2–50", "Synthetic cells"), ("0.97", "AUC at N=8"),
                                   ("0.46", "Real-data AUC"), ("3", "Code bugs fixed")]):
    add_stat_box(slide, 0.8 + i * 2.9, 4.5, num, label)

# ── SLIDE 2: Why This Matters ───────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, 0.8, 0.5, 11.5, 0.6, "Why this matters", font_size=30, bold=True, color=WHITE)
add_text_box(slide, 0.8, 1.2, 11.5, 0.5,
             "Grid batteries face a critical decision: accept or reject each service request.",
             font_size=18, color=GRAY)
add_bullet_frame(slide, 0.8, 1.9, 5.5, 2.0, [
    "✅ Say yes when healthy → earn revenue",
    "❌ Say yes when failing → blackout, penalties",
    "❌ Say no when fine → lost revenue",
], font_size=16)
add_text_box(slide, 0.8, 3.8, 11.5, 0.5,
             "Shikdar & Laaksonen (2026) built an ML framework to predict failure probability across 10–50 cycle horizons.",
             font_size=16, color=GRAY)
add_bullet_frame(slide, 0.8, 4.5, 11.5, 1.5, [
    "Reported AUC = 0.944 on 37 batteries",
    "Failure rate reduction: 10.3% → 2.95%",
    "Can we reproduce this with the 4-cell NASA dataset?",
], font_size=16, color=YELLOW)

# ── SLIDE 3: Reproducibility Pipeline ───────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, 0.8, 0.5, 11.5, 0.6, "Our Reproducibility Pipeline", font_size=30, bold=True, color=WHITE)
add_bullet_frame(slide, 0.8, 1.3, 11.5, 5.0, [
    "Data layer: NASALoader parses .mat files → normalized DataFrame with SOH, voltage, current, temperature features",
    "Model layer: XGBoost multi-output classifier per horizon, isotonic calibration on held-out validation fold",
    "Dispatch layer: Threshold policies, continuous derating, Monte Carlo market simulator (AR(1) price model)",
    "Evaluation: Leave-one-battery-out CV, AUC/Brier/ECE metrics, per-fold variance reporting",
    "",
    "Run with:  python run_all.py --quick  (completes in < 10 seconds)",
], font_size=16)

# ── SLIDE 4: The Data ───────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, 0.8, 0.5, 11.5, 0.6, "The NASA 4-Cell Dataset", font_size=30, bold=True, color=WHITE)

add_stat_box(slide, 0.8, 1.3, "B0005", "168 cycles, min SOH=0.699")
add_stat_box(slide, 3.7, 1.3, "B0006", "168 cycles, min SOH=0.572")
add_stat_box(slide, 6.6, 1.3, "B0007", "168 cycles, min SOH=0.744")
add_stat_box(slide, 9.5, 1.3, "B0018", "132 cycles, min SOH=0.729")

add_text_box(slide, 0.8, 3.5, 11.5, 0.5,
             "Key limitation:", font_size=18, bold=True, color=YELLOW)
add_bullet_frame(slide, 0.8, 4.0, 11.5, 2.5, [
    "Only 64 cycles fall below the EOL threshold (SOH < 0.70)",
    "62 of those come from a single cell (B0006)",
    "Leave-battery-out CV: train on 3 cells, test on 1",
    "Most folds have zero positive samples → AUC cannot be computed",
], font_size=16)

# ── SLIDE 5: Main Results ───────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, 0.8, 0.5, 11.5, 0.6, "Results: Near-Random Discrimination", font_size=30, bold=True, color=WHITE)
add_table(slide, 0.8, 1.3, 11.5,
          ["Horizon", "Raw AUC", "Calibrated AUC", "Per-fold (mean±std)"],
          [
              ["H=10", "0.26", "0.50", "— (2/4 folds NaN)"],
              ["H=20", "0.26", "0.50", "— (2/4 folds NaN)"],
              ["H=30", "0.60", "0.50", "0.50±0.00"],
              ["H=50", "0.69", "0.50", "0.50±0.00"],
              ["Macro avg", "0.46", "0.50", "0.50±0.00"],
          ])
add_text_box(slide, 0.8, 4.2, 11.5, 0.5,
             "Random guessing = AUC 0.50. The raw model produces near-constant predictions (no rank variation). Isotonic regression preserves rank order — with no rank order to preserve, AUC is undefined and conventionally reported as 0.5.",
             font_size=14, color=RED)
add_text_box(slide, 0.8, 4.8, 11.5, 0.5,
             "Original paper achieved AUC 0.944 on 37 batteries — the framework needs substantial data.",
             font_size=14, color=GRAY)

# ── SLIDE 6: Scaling Study ──────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, 0.8, 0.5, 11.5, 0.6, "Synthetic Scaling Study: AUC vs Dataset Size", font_size=30, bold=True, color=WHITE)
add_text_box(slide, 0.8, 1.2, 11.5, 0.5,
             "For N in {2, 3, 5, 8, 12, 20, 30, 50} with 3 Monte Carlo seeds each:",
             font_size=16, color=GRAY)
add_table(slide, 0.8, 1.9, 11.5,
          ["N cells", "Mean AUC", "±1 Std", "Regime"],
          [
              ["2", "0.84", "0.009", "Insufficient"],
              ["3", "0.90", "0.019", "Insufficient"],
              ["5", "0.93", "0.007", "Insufficient"],
              ["8", "0.97", "0.002", "Marginal"],
              ["12", "0.98", "0.002", "Reliable"],
              ["20", "0.99", "0.001", "Reliable"],
              ["30", "0.99", "0.001", "Reliable"],
              ["50", "0.99", "0.001", "Reliable"],
          ])
add_text_box(slide, 0.8, 4.6, 11.5, 0.5,
             "Diminishing returns beyond N=8–12. Real-world data will need more (est. 15–25 for AUC > 0.95).",
             font_size=16, color=YELLOW)
add_text_box(slide, 0.8, 5.2, 11.5, 0.4,
             "Per-horizon breakdown (representative seed):  N=3 → H=10:0.88, H=20:0.89, H=30:0.91, H=50:0.92 — all horizons improve together.",
             font_size=13, color=GRAY)
add_text_box(slide, 0.8, 5.7, 11.5, 0.4,
             "Diversity caveat: Synthetic data maximizes cell diversity. Real batteries share chemistry/manufacturer — correlated degradation reduces effective diversity.",
             font_size=13, color=GRAY)
add_text_box(slide, 0.8, 6.1, 11.5, 0.4,
             "Variance note: N ≤ 3 estimates based on 2–3 CV folds, interpret cautiously.  |  Comput. cost: full scaling ~2–3 h, NASA real-data < 10 s.",
             font_size=13, color=GRAY)

# ── SLIDE 7: Three Code Bugs ────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, 0.8, 0.5, 11.5, 0.6, "Three Hidden Bugs We Found", font_size=30, bold=True, color=WHITE)
add_card(slide, 0.8, 1.3, 11.5, 1.6,
         "Bug 1: Calibration Data Leakage",
         "Original code fit the isotonic calibrator on test data — then evaluated on that same data. This inflated AUC from 0.46 to an apparent 0.74. Fix: fit on a held-out validation fold.",
         title_color=RED)
add_card(slide, 0.8, 3.1, 11.5, 1.6,
         "Bug 2: Energy Unit Confusion",
         "Revenue = energy (kWh) × price ($/MWh) — without dividing by 1000, revenue was overstated by 1000× ($3,780 vs $3.78).",
         title_color=RED)
add_card(slide, 0.8, 4.9, 11.5, 1.6,
         "Bug 3: Inconsistent Baseline Metric",
         "Baseline failure rate used label density across all horizons; model rows used conditional dispatch-based metric. These are incomparable.",
         title_color=RED)

# ── SLIDE 8: Ablation ───────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, 0.8, 0.5, 11.5, 0.6, "Ablation Study", font_size=30, bold=True, color=WHITE)
add_table(slide, 0.8, 1.3, 11.5,
          ["Configuration", "Failure Rate", "Energy"],
          [
              ["Always dispatch", "0.0063", "318.0"],
              ["+ Raw hazard model", "0.0063", "318.0"],
              ["+ Probability calibration", "0.0063", "318.0"],
              ["+ Composite failure labels", "0.0110", "318.0"],
          ])
add_text_box(slide, 0.8, 4.0, 11.5, 1.0,
             "Composite labels increase the failure rate — the NASA dataset lacks resistance measurements,\nso composite labels add noise without predictive signal.",
             font_size=16, color=RED)

# ── SLIDE 9: Takeaways ──────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, 0.8, 0.5, 11.5, 0.6, "Key Takeaways", font_size=30, bold=True, color=WHITE)

add_card(slide, 0.8, 1.3, 11.5, 1.5,
         "First quantitative scaling curve for battery ML",
         "AUC vs N: N ≤ 5 = insufficient, 5 < N < 12 = marginal, N ≥ 12 = reliable. Real data estimated at 15–25+ for AUC > 0.95.",
         title_color=ACCENT)
add_card(slide, 0.8, 3.0, 11.5, 1.5,
         "Open-source framework with correct methodology",
         "Three bugs documented and fixed: calibration leakage (+0.28 AUC), energy unit error (1000×), inconsistent baselines. All code public.",
         title_color=ACCENT)
add_card(slide, 0.8, 4.7, 11.5, 1.5,
         "Honest real-data validation",
         "On 4 NASA cells, AUC = 0.46 — confirming the framework needs more data than one public benchmark. Anyone can reproduce.",
         title_color=ACCENT)

# ── SLIDE 10: Software ──────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, 0.8, 0.5, 11.5, 0.6, "Our Python Package", font_size=30, bold=True, color=WHITE)
add_bullet_frame(slide, 0.8, 1.3, 11.5, 4.5, [
    "Modular — data, models, dispatch, evaluation layers",
    "Correct calibration — fit on validation fold, not test set",
    "Unit-aware economics — correct kWh ↔ $/MWh conversion",
    "Consistent metrics — same dispatch-based metric for all ablation rows",
    "Scaling experiment — --expt scaling runs N=2..50 Monte Carlo study",
    "Synthetic data fallback — pipeline runs without external files",
    "Full configuration — single config.yaml controls everything",
    "",
    "Code: github.com/teamdynamic/battery-reliability-extension",
    "DOI: 10.5281/zenodo.XXXXXXX",
], font_size=16)

# ── SLIDE 11: Limitations & Future Work ─────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, 0.8, 0.5, 11.5, 0.6, "Limitations & Future Work", font_size=30, bold=True, color=WHITE)
add_text_box(slide, 0.8, 1.2, 5.5, 0.4, "Current limitations", font_size=20, bold=True, color=RED)
add_bullet_frame(slide, 0.8, 1.7, 5.5, 3.5, [
    "Only 4 cells, all 18,650 format under lab conditions",
    "No deep learning (LSTM, TCN, Transformer) — CPU constraints",
    "Simplified AR(1) price model",
    "CALCE cross-chemistry evaluation incomplete",
], font_size=14)
add_text_box(slide, 6.8, 1.2, 5.5, 0.4, "Future work", font_size=20, bold=True, color=GREEN)
add_bullet_frame(slide, 6.8, 1.7, 5.5, 3.5, [
    "Evaluate on larger public datasets (10+ batteries)",
    "Cross-chemistry validation (LFP, NMC)",
    "GPU-based deep learning benchmarks",
    "Real price data instead of AR(1) simulation",
], font_size=14)

# ── SLIDE 12: Thank You ─────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, 0.8, 2.0, 11.5, 1.0, "Thank You", font_size=44, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 3.2, 11.5, 0.6, "Questions?  Feedback?",
             font_size=22, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 4.5, 11.5, 0.5, "Paper:  Extension_Paper.docx",
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 5.0, 11.5, 0.5, "Code:  github.com/teamdynamic/battery-reliability-extension",
             font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 6.0, 11.5, 0.4, "Team Dynamic  •  June 2026",
             font_size=14, color=DIM, alignment=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Size: {os.path.getsize(out_path) // 1024} KB")
